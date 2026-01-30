from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

class SlidesGenerator:
    def __init__(self, output_dir="exports"):
        self.output_dir = output_dir
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

    def generate_slides(self, data: dict, filename: str):
        prs = Presentation()
        
        # Style Dark Mode de base (fond noir)
        self._apply_dark_theme(prs)

        # 1. Slide de Titre
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = data.get("title", "Rapport DocTerra").upper()
        subtitle.text = f"Synthèse Stratégique Alchimique\nGénéré le {data.get('date', '30/01/2026')}"
        
        self._style_text(title, size=Pt(44), bold=True)
        self._style_text(subtitle, size=Pt(20), color=RGBColor(0, 209, 255))

        # 2. Slides de Contenu (une par section)
        for section in data.get("sections", []):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            # Titre de la slide
            slide_title = slide.shapes.title
            slide_title.text = section.get("title", "Section").upper()
            self._style_text(slide_title, size=Pt(32), bold=True, color=RGBColor(0, 209, 255))
            
            # Contenu (bullet points)
            content = slide.placeholders[1]
            content.text = section.get("content", section.get("brief", ""))
            self._style_text(content, size=Pt(18))

        # Sauvegarde
        if not filename.endswith(".pptx"):
            filename += ".pptx"
        output_path = os.path.join(self.output_dir, filename)
        prs.save(output_path)
        return output_path

    def _apply_dark_theme(self, prs):
        # Pour une version plus complexe, on utiliserait un template .potx
        # Ici on simplifie en changeant les couleurs par défaut via l'objet slide si besoin
        pass

    def _style_text(self, shape, size=None, bold=False, color=RGBColor(255, 255, 255)):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if size:
                    run.font.size = size
                run.font.bold = bold
                run.font.color.rgb = color
                run.font.name = "Arial"
